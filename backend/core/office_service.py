"""
Office Automation Service for Atom

Provides core utilities for reading, modifying, and rendering Word (.docx),
Excel (.xlsx), and PowerPoint (.pptx) documents without native Office dependencies.
"""

import logging
import os
import re
import io
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple, Union

# Third-party imports (lazy loaded / imported defensively)
import openpyxl
import docx
try:
    import pptx
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import mammoth
    MAMMOTH_AVAILABLE = True
except ImportError:
    MAMMOTH_AVAILABLE = False

try:
    import xlsx2html
    XLSX2HTML_AVAILABLE = True
except ImportError:
    XLSX2HTML_AVAILABLE = False

logger = logging.getLogger(__name__)


def _validate_office_path(file_path: str) -> str:
    """Contain user-supplied file paths to the configured office directory.

    Office endpoints accept a user-controlled ``file_path``; without containment
    that's arbitrary file read/write (path traversal — reading /etc/passwd or
    overwriting config). Resolve to an absolute path and require it to live
    under ATOM_OFFICE_DIR (default ./data/office). Raises ValueError otherwise.
    """
    if not file_path:
        raise ValueError("file_path is required")
    base = Path(os.getenv("ATOM_OFFICE_DIR", os.path.join("data", "office"))).resolve()
    try:
        resolved = Path(file_path).resolve()
    except (OSError, ValueError) as e:
        raise ValueError("Invalid file path")
    # Resolve symlinks/relative escapes: the resolved path must be under base.
    if base not in resolved.parents and resolved != base:
        raise ValueError(
            f"Access denied: path '{file_path}' is outside the allowed office directory"
        )
    return str(resolved)


class ExcelManager:
    """Manages Excel sheet operations using openpyxl."""

    @staticmethod
    def parse_path(path: str) -> Tuple[str, str]:
        """
        Parses a DOM-like path.
        Example: '/Sheet1/A1:B10' -> ('Sheet1', 'A1:B10')
                 '/Sheet1/A1' -> ('Sheet1', 'A1')
                 'A1' -> ('', 'A1')
        """
        path = path.strip()
        if path.startswith('/'):
            parts = [p for p in path.split('/') if p]
            if len(parts) >= 2:
                return parts[0], '/'.join(parts[1:])
            elif len(parts) == 1:
                return parts[0], ""
        return "", path

    def read_range(self, file_path: str, cell_path: str) -> Dict[str, Any]:
        """Read values from a cell or cell range."""
        try:
            file_path = _validate_office_path(file_path)
        except ValueError as e:
            return {"success": False, "error": str(e)}
        if not os.path.exists(file_path):
            return {"success": False, "error": f"File not found"}

        try:
            wb = openpyxl.load_workbook(file_path, data_only=True)
            sheet_name, coordinate = self.parse_path(cell_path)

            if sheet_name and sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
            else:
                ws = wb.active
                sheet_name = ws.title

            if not coordinate:
                # Return sheet overview or all cells if range is empty
                return {
                    "success": True,
                    "sheet_name": sheet_name,
                    "dimensions": ws.dimensions,
                    "sheet_names": wb.sheetnames
                }

            # Check if it's a range or a single cell
            if ":" in coordinate:
                cells_range = ws[coordinate]
                data = []
                for row in cells_range:
                    row_data = []
                    for cell in row:
                        row_data.append({
                            "cell_ref": cell.coordinate,
                            "value": cell.value,
                            "cell_type": "formula" if cell.value and str(cell.value).startswith('=') else "text"
                        })
                    data.append(row_data)
                return {
                    "success": True,
                    "sheet_name": sheet_name,
                    "coordinate": coordinate,
                    "cells": data
                }
            else:
                cell = ws[coordinate]
                raw_wb = openpyxl.load_workbook(file_path, data_only=False)
                raw_ws = raw_wb[sheet_name] if sheet_name in raw_wb.sheetnames else raw_wb.active
                raw_cell = raw_ws[coordinate]
                
                is_formula = str(raw_cell.value).startswith('=') if raw_cell.value else False
                formula = str(raw_cell.value) if is_formula else None

                return {
                    "success": True,
                    "sheet_name": sheet_name,
                    "coordinate": coordinate,
                    "value": cell.value,
                    "formula": formula,
                    "cell_type": "formula" if is_formula else "text"
                }
        except Exception as e:
            logger.error(f"Error reading Excel range: {e}")
            return {"success": False, "error": "Failed to read Excel range"}

    def create_spreadsheet(self, file_path: str, rows: List[List[Any]]) -> Dict[str, Any]:
        """Create a new .xlsx from row data (chat-table drafts → office
        canvas). One workbook open/save, unlike per-cell write_cell.
        Numeric-looking strings become real numbers so Excel formulas and
        sums work instead of tripping over text cells."""
        def coerce(value: Any) -> Any:
            if isinstance(value, str):
                stripped = value.strip()
                if re.fullmatch(r"-?\d+", stripped):
                    return int(stripped)
                if re.fullmatch(r"-?\d*\.\d+", stripped):
                    return float(stripped)
            return value

        try:
            file_path = _validate_office_path(file_path)
        except ValueError as e:
            return {"success": False, "error": str(e)}
        try:
            wb = openpyxl.Workbook()
            ws = wb.active
            ws.title = "Sheet1"
            for row in rows or []:
                ws.append([coerce(cell) for cell in row])
            os.makedirs(os.path.dirname(file_path) or ".", exist_ok=True)
            wb.save(file_path)
            return {"success": True, "file_path": file_path, "rows": len(rows or [])}
        except Exception as e:
            return {"success": False, "error": str(e)}

    def write_cell(self, file_path: str, cell_path: str, value: Any, is_formula: bool = False) -> Dict[str, Any]:
        """Write value or formula to a cell."""
        try:
            file_path = _validate_office_path(file_path)
        except ValueError as e:
            return {"success": False, "error": str(e)}
        try:
            if os.path.exists(file_path):
                wb = openpyxl.load_workbook(file_path)
            else:
                wb = openpyxl.Workbook()
            sheet_name, coordinate = self.parse_path(cell_path)

            if sheet_name:
                if sheet_name not in wb.sheetnames:
                    wb.create_sheet(title=sheet_name)
                ws = wb[sheet_name]
            else:
                ws = wb.active
                sheet_name = ws.title

            if not coordinate:
                return {"success": False, "error": "Cell coordinate not specified"}

            # Cast value appropriately if not a formula
            if is_formula and not str(value).startswith('='):
                value = f"={value}"
            elif not is_formula:
                # Try to cast to float/int if possible
                if isinstance(value, str):
                    if value.isdigit():
                        value = int(value)
                    else:
                        try:
                            value = float(value)
                        except ValueError:
                            pass

            ws[coordinate] = value
            if not is_formula and isinstance(value, str) and value.startswith("="):
                cell = ws[coordinate]
                cell.data_type = "s"
            wb.save(file_path)

            # Recalculate via the workbook runtime so the agent sees computed
            # results immediately (not stale formula strings). Best-effort:
            # a runtime failure doesn't break the write.
            computed_value = value
            try:
                from core.workbook_runtime import get_workbook_runtime
                runtime = get_workbook_runtime()
                if runtime.can_evaluate and is_formula:
                    # Run recalc synchronously (write_cell is a sync method).
                    import asyncio
                    try:
                        loop = asyncio.new_event_loop()
                        loop.run_until_complete(runtime.recalculate(file_path))
                        loop.close()
                    except Exception:
                        pass  # Recalc failed — keep the formula as the value.
                    # Read the now-computed value.
                    wb2 = openpyxl.load_workbook(file_path, data_only=True)
                    ws2 = wb2[sheet_name] if sheet_name in wb2.sheetnames else wb2.active
                    computed_value = ws2[coordinate].value
            except Exception as recalc_err:
                logger.debug(f"Recalc after write skipped (non-fatal): {recalc_err}")

            return {
                "success": True,
                "sheet_name": sheet_name,
                "coordinate": coordinate,
                "value": computed_value,
                "formula": value if is_formula else None,
                "message": f"Updated {sheet_name}!{coordinate} successfully"
            }
        except Exception as e:
            logger.error(f"Error writing Excel cell: {e}")
            return {"success": False, "error": "Failed to write Excel cell"}

    @staticmethod
    async def insert_rows(file_path: str, sheet_name: str, row: int, count: int = 1) -> Dict[str, Any]:
        """Insert rows and recalculate formulas to maintain references."""
        from core.workbook_runtime import get_workbook_runtime
        try:
            file_path = _validate_office_path(file_path)
        except ValueError as e:
            return {"success": False, "error": str(e)}
        if not os.path.exists(file_path):
            return {"success": False, "error": "File not found"}
        runtime = get_workbook_runtime()
        return await runtime.insert_rows(file_path, sheet_name, row, count)

    @staticmethod
    async def insert_columns(file_path: str, sheet_name: str, col: int, count: int = 1) -> Dict[str, Any]:
        """Insert columns and recalculate formulas to maintain references."""
        from core.workbook_runtime import get_workbook_runtime
        try:
            file_path = _validate_office_path(file_path)
        except ValueError as e:
            return {"success": False, "error": str(e)}
        if not os.path.exists(file_path):
            return {"success": False, "error": "File not found"}
        runtime = get_workbook_runtime()
        return await runtime.insert_cols(file_path, sheet_name, col, count)

    @staticmethod
    async def get_evaluated_range(file_path: str, cell_path: str) -> Dict[str, Any]:
        """Read a range with freshly evaluated formula results."""
        from core.workbook_runtime import get_workbook_runtime
        try:
            file_path = _validate_office_path(file_path)
        except ValueError as e:
            return {"success": False, "error": str(e)}
        if not os.path.exists(file_path):
            return {"success": False, "error": "File not found"}
        runtime = get_workbook_runtime()
        sheet_name, cell_range = ExcelManager.parse_path(cell_path)
        parts = cell_range.split(":") if ":" in cell_range else [cell_range]
        start = parts[0] if parts else "A1"
        end = parts[1] if len(parts) > 1 else None
        return await runtime.get_evaluated_range(file_path, sheet_name, start, end)

    @staticmethod
    async def recalculate(file_path: str) -> Dict[str, Any]:
        """Force recalculation of all formulas in the workbook."""
        from core.workbook_runtime import get_workbook_runtime
        from pathlib import Path
        try:
            file_path = _validate_office_path(file_path)
        except ValueError as e:
            return {"success": False, "error": str(e)}
        if not os.path.exists(file_path):
            return {"success": False, "error": "File not found"}
        runtime = get_workbook_runtime()
        await runtime.recalculate(Path(file_path))
        return {"success": True, "engine": runtime.engine}

    @staticmethod
    async def add_pivot_table(
        file_path: str, sheet_name: str, pivot_sheet_name: str,
        data_range: str, rows: List[str], columns: List[str],
        values: List[Dict[str, str]]
    ) -> Dict[str, Any]:
        """Create a styled pivot table summary sheet."""
        from core.workbook_runtime import get_workbook_runtime
        try:
            file_path = _validate_office_path(file_path)
        except ValueError as e:
            return {"success": False, "error": str(e)}
        if not os.path.exists(file_path):
            return {"success": False, "error": "File not found"}
        runtime = get_workbook_runtime()
        return await runtime.add_pivot_table(
            file_path, sheet_name, pivot_sheet_name, data_range, rows, columns, values
        )

    @staticmethod
    async def run_excel_macro(file_path: str, macro_name: str) -> Dict[str, Any]:
        """Run macro inside workbook via sandboxed execution."""
        from core.workbook_runtime import get_workbook_runtime
        try:
            file_path = _validate_office_path(file_path)
        except ValueError as e:
            return {"success": False, "error": str(e)}
        if not os.path.exists(file_path):
            return {"success": False, "error": "File not found"}
        runtime = get_workbook_runtime()
        return await runtime.run_macro(file_path, macro_name)


class WordManager:
    """Manages Word document operations using python-docx."""

    def read_document(self, file_path: str) -> Dict[str, Any]:
        """Read content from a Word document."""
        try:
            file_path = _validate_office_path(file_path)
        except ValueError as e:
            return {"success": False, "error": str(e)}
        if not os.path.exists(file_path):
            return {"success": False, "error": "File not found"}

        try:
            doc = docx.Document(file_path)
            paragraphs = []
            for i, p in enumerate(doc.paragraphs):
                if p.text.strip():
                    paragraphs.append({
                        "index": i,
                        "text": p.text,
                        "style": p.style.name
                    })

            tables = []
            for t_idx, table in enumerate(doc.tables):
                table_data = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(row_data)
                tables.append({
                    "index": t_idx,
                    "rows": table_data
                })

            return {
                "success": True,
                "paragraphs": paragraphs,
                "tables": tables,
                "metadata": {
                    "paragraphs_count": len(doc.paragraphs),
                    "tables_count": len(doc.tables)
                }
            }
        except Exception as e:
            logger.error(f"Error reading Word document: {e}")
            return {"success": False, "error": "Failed to read Word document"}

    def modify_document(self, file_path: str, action: str, content: str, options: dict = None) -> Dict[str, Any]:
        """Modify a Word document."""
        options = options or {}
        try:
            file_path = _validate_office_path(file_path)
        except ValueError as e:
            return {"success": False, "error": str(e)}
        try:
            doc = docx.Document(file_path) if os.path.exists(file_path) else docx.Document()

            if action == "append":
                style = options.get("style", "Normal")
                doc.add_paragraph(content, style=style)
            elif action == "replace":
                target = options.get("target")
                if not target:
                    return {"success": False, "error": "Replace action requires a target placeholder"}
                
                replaced_count = 0
                for p in doc.paragraphs:
                    if target in p.text:
                        p.text = p.text.replace(target, content)
                        replaced_count += 1
                
                # Check tables too
                for table in doc.tables:
                    for row in table.rows:
                        for cell in row.cells:
                            if target in cell.text:
                                cell.text = cell.text.replace(target, content)
                                replaced_count += 1
            else:
                return {"success": False, "error": f"Unknown modification action: {action}"}

            doc.save(file_path)
            return {
                "success": True,
                "message": f"Word document modified successfully (Action: {action})"
            }
        except Exception as e:
            logger.error(f"Error modifying Word document: {e}")
            return {"success": False, "error": "Failed to modify Word document"}


class PowerPointManager:
    """Manages PowerPoint slide decks using python-pptx."""

    def read_slides(self, file_path: str) -> Dict[str, Any]:
        """Read content from a PowerPoint deck."""
        if not PPTX_AVAILABLE:
            return {"success": False, "error": "python-pptx library not installed"}
        try:
            file_path = _validate_office_path(file_path)
        except ValueError as e:
            return {"success": False, "error": str(e)}
        if not os.path.exists(file_path):
            return {"success": False, "error": "File not found"}

        try:
            prs = pptx.Presentation(file_path)
            slides = []

            for idx, slide in enumerate(prs.slides):
                shapes = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        shapes.append({
                            "type": "text",
                            "name": shape.name,
                            "text": shape.text_frame.text
                        })
                    elif shape.has_table:
                        table_data = []
                        for row in shape.table.rows:
                            row_data = [cell.text for cell in row.cells]
                            table_data.append(row_data)
                        shapes.append({
                            "type": "table",
                            "name": shape.name,
                            "table": table_data
                        })

                slides.append({
                    "slide_index": idx,
                    "shapes": shapes
                })

            return {
                "success": True,
                "slides": slides,
                "slide_count": len(prs.slides)
            }
        except Exception as e:
            logger.error(f"Error reading PowerPoint: {e}")
            return {"success": False, "error": "Failed to read PowerPoint"}

    def modify_slides(self, file_path: str, action: str, options: dict) -> Dict[str, Any]:
        """Modify slide deck."""
        if not PPTX_AVAILABLE:
            return {"success": False, "error": "python-pptx library not installed"}
        try:
            file_path = _validate_office_path(file_path)
        except ValueError as e:
            return {"success": False, "error": str(e)}
        try:
            prs = pptx.Presentation(file_path) if os.path.exists(file_path) else pptx.Presentation()

            if action == "add_slide":
                # Standard slide layouts: 0=Title, 1=Title+Content, etc.
                layout_idx = options.get("layout_idx", 1)
                if layout_idx >= len(prs.slide_layouts):
                    layout_idx = 1
                
                blank_layout = prs.slide_layouts[layout_idx]
                slide = prs.slides.add_slide(blank_layout)
                
                title = options.get("title")
                if title and slide.shapes.title:
                    slide.shapes.title.text = title
                
                content = options.get("content")
                if content and len(slide.placeholders) > 1:
                    slide.placeholders[1].text = content
            elif action == "update_slide":
                # Canvas co-editing: update title/body text of an existing slide
                # (1-based slide_number) without touching other slides.
                slide_number = int(options.get("slide_number", 0))
                if slide_number < 1 or slide_number > len(prs.slides):
                    return {"success": False, "error": f"Slide number out of range: {slide_number}"}
                slide = prs.slides[slide_number - 1]

                new_title = options.get("title")
                if new_title is not None and slide.shapes.title is not None:
                    slide.shapes.title.text = str(new_title)

                new_content = options.get("content")
                if new_content is not None:
                    target = None
                    for ph in slide.placeholders:
                        try:
                            idx = ph.placeholder_format.idx
                        except Exception:
                            continue
                        if idx == 1 and ph.has_text_frame:
                            target = ph
                            break
                    if target is None:
                        title_shape = slide.shapes.title
                        for shape in slide.shapes:
                            if shape.has_text_frame and shape is not title_shape:
                                target = shape
                                break
                    if target is not None:
                        target.text_frame.text = str(new_content)

            elif action == "set_shape_text":
                # Canvas co-edit: replace the text of an existing shape while
                # preserving the shape itself, its position, and its layout.
                slide_index = options.get("slide_index")
                shape_name = options.get("shape_name")
                text = options.get("text", "")
                if not isinstance(slide_index, int) or slide_index < 0 or slide_index >= len(prs.slides):
                    return {"success": False, "error": f"Invalid slide_index: {slide_index}"}
                if not shape_name:
                    return {"success": False, "error": "shape_name required for set_shape_text"}

                slide = prs.slides[slide_index]
                shape = next((s for s in slide.shapes if s.name == shape_name), None)
                if shape is None or not shape.has_text_frame:
                    return {"success": False, "error": f"Text shape '{shape_name}' not found on slide {slide_index}"}

                # Setting .text replaces all runs with a single run inheriting
                # the paragraph's style — the least destructive simple edit.
                shape.text_frame.text = str(text)
            else:
                return {"success": False, "error": f"Unknown PowerPoint action: {action}"}

            prs.save(file_path)
            return {"success": True, "message": f"PowerPoint modified successfully (Action: {action})"}
        except Exception as e:
            logger.error(f"Error modifying PowerPoint: {e}")
            return {"success": False, "error": "Failed to modify PowerPoint"}


class DocumentRenderer:
    """Renders Word, Excel, and PowerPoint documents to HTML or PNG images."""

    @staticmethod
    def render_to_html(file_path: str) -> Dict[str, Any]:
        """Convert a document into an HTML string for previewing."""
        try:
            file_path = _validate_office_path(file_path)
        except ValueError as e:
            return {"success": False, "error": str(e)}
        ext = Path(file_path).suffix.lower()

        if ext == ".docx":
            if not MAMMOTH_AVAILABLE:
                return {"success": False, "error": "mammoth library not installed"}
            try:
                with open(file_path, "rb") as docx_file:
                    result = mammoth.convert_to_html(docx_file)
                    # Mammoth provides conversion messages
                    return {
                        "success": True,
                        "html": f"<div class='office-word-preview'>{result.value}</div>",
                        "warnings": [m.message for m in result.messages]
                    }
            except Exception as e:
                logger.error(f"Error rendering Word to HTML: {e}")
                return {"success": False, "error": "Failed rendering Word to HTML"}

        elif ext == ".xlsx":
            # Use the workbook runtime for pixel-accurate rendering (LibreOffice
            # when available — includes conditional formatting, charts, and
            # evaluated formulas). Falls back to a basic openpyxl HTML table.
            try:
                from core.workbook_runtime import get_workbook_runtime
                import asyncio as _asyncio
                runtime = get_workbook_runtime()
                try:
                    loop = _asyncio.get_event_loop()
                    if loop.is_running():
                        # We're in an async context — can't await here directly.
                        # Fall back to basic render (the sync path).
                        html_val = runtime._render_html_basic(Path(file_path))
                    else:
                        html_val = loop.run_until_complete(runtime.render_to_html(file_path))
                except RuntimeError:
                    html_val = runtime._render_html_basic(Path(file_path))
                return {
                    "success": True,
                    "html": f"<div class='office-excel-preview'>{html_val}</div>",
                    "engine": runtime.engine,
                }
            except Exception as e:
                logger.error(f"Error rendering Excel to HTML: {e}")
                return {"success": False, "error": "Failed rendering Excel to HTML"}

        elif ext == ".pptx":
            if not PPTX_AVAILABLE:
                return {"success": False, "error": "python-pptx library not installed"}
            try:
                prs = pptx.Presentation(file_path)
                html_slides = []
                for i, slide in enumerate(prs.slides):
                    texts = []
                    for shape in slide.shapes:
                        if shape.has_text_frame and shape.text_frame.text:
                            texts.append(f"<p>{shape.text_frame.text}</p>")
                    
                    slide_content = "\n".join(texts)
                    html_slides.append(f"""
                        <div class="slide" style="border:1px solid #ccc; padding:20px; margin-bottom:20px; aspect-ratio:16/9; background:#fff; box-shadow: 0 4px 6px rgba(0,0,0,0.1);">
                            <h3>Slide {i+1}</h3>
                            <div class="slide-content">{slide_content}</div>
                        </div>
                    """)
                return {
                    "success": True,
                    "html": f"<div class='office-pptx-preview' style='background:#f4f4f4; padding:20px;'>{''.join(html_slides)}</div>"
                }
            except Exception as e:
                logger.error(f"Error rendering PPTX to HTML: {e}")
                return {"success": False, "error": "Failed rendering PPTX to HTML"}

        return {"success": False, "error": f"Unsupported format: {ext}"}


class OfficeService:
    """Primary service coordinating Office document manipulation."""

    def __init__(self):
        self.excel = ExcelManager()
        self.word = WordManager()
        self.pptx = PowerPointManager()
        self.renderer = DocumentRenderer()

    def get_manager_for_file(self, file_path: str) -> Any:
        """Get the appropriate manager depending on file suffix."""
        ext = Path(file_path).suffix.lower()
        if ext == ".xlsx":
            return self.excel
        elif ext == ".docx":
            return self.word
        elif ext == ".pptx":
            return self.pptx
        else:
            raise ValueError(f"Unsupported file format: {ext}")
