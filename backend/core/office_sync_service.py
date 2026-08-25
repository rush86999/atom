"""
Office Synchronization Service for Atom

Bridges the local filesystem documents (.docx, .xlsx, .pptx) and the active Canvas state.
Ensures edits made by the user on the Canvas are saved to disk, and edits made by the agent
on disk are pushed to the Canvas.

In addition to file/Canvas sync, every meaningful document change is ingested into Atom's
memory (knowledge graph + LanceDB) so the agent "remembers" the content of quotes, POs,
price lists, and invoices it generates or edits.

Co-editing model (#39):
  - User edits on the canvas POST /api/v1/office/sync-update → file is modified in place
    (docx paragraph sync preserves styles/tables/images; xlsx cell writes recalc formulas;
    pptx slide text updates round-trip).
  - Every file change broadcasts a STRUCTURED content snapshot (grid cells / paragraphs /
    slides) plus rendered HTML over WebSocket on BOTH `canvas:{id}` and `user:{uid}`
    channels so any open canvas — chat panel or /canvas/{id} workspace — re-renders.
  - Agent writes through tools/office_tool.py call notify_file_canvases() so bound
    canvases pick up agent-side edits too.
"""

import asyncio
import logging
import os
import uuid
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

import openpyxl
import docx  # python-docx — used for Word-doc canvas→file sync
from sqlalchemy.orm import Session

from core.office_service import OfficeService, _validate_office_path
from core.models import Canvas, CanvasAudit
from core.websockets import manager as ws_manager

try:
    import pptx
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

logger = logging.getLogger(__name__)

# Extension → (WS component name the frontend renders, DB canvas_type).
OFFICE_COMPONENT_MAP: Dict[str, Tuple[str, str]] = {
    ".xlsx": ("office_excel", "sheets"),
    ".docx": ("office_word", "docs"),
    ".pptx": ("office_pptx", "presentation"),
}


class OfficeSyncService:
    """Service coordinates bi-directional sync between files and canvas UI."""

    def __init__(self, db: Session):
        self.db = db
        self.office = OfficeService()

    # ------------------------------------------------------------------
    # Structured content snapshots (file → UI)
    # ------------------------------------------------------------------

    @staticmethod
    def _read_structured_content(file_path: str) -> Dict[str, Any]:
        """Read an office file into the structured shape the editable canvas
        components render. Best-effort: raises nothing, returns {} on failure.

        - xlsx → {format, active_sheet, sheet_names, sheets:[{name, rows}], formulas}
          rows hold computed values; formulas maps Sheet→{coord:'=...'} for the
          formula bar.
        - docx → {format, text, paragraphs:[{index,text,style}]} where text is
          ALL paragraphs joined by newline (1:1 mapping for in-place sync back).
        - pptx → {format, slides:[{slide_number,title,content}]}
        """
        try:
            ext = Path(file_path).suffix.lower()
            if ext == ".xlsx":
                wb = openpyxl.load_workbook(file_path, data_only=True)
                wb_raw = openpyxl.load_workbook(file_path, data_only=False)
                sheets: List[Dict[str, Any]] = []
                formulas: Dict[str, Dict[str, str]] = {}
                for ws in wb.worksheets:
                    rows: List[List[Any]] = []
                    for row in ws.iter_rows():
                        rows.append([c.value for c in row])
                    sheets.append({"name": ws.title, "rows": rows})
                    # data_only=True strips formulas — read raw for the bar.
                    fmap: Dict[str, str] = {}
                    for rrow in wb_raw[ws.title].iter_rows():
                        for c in rrow:
                            if isinstance(c.value, str) and c.value.startswith("="):
                                fmap[c.coordinate] = c.value
                    if fmap:
                        formulas[ws.title] = fmap
                return {
                    "format": "xlsx",
                    "active_sheet": wb.active.title,
                    "sheet_names": wb.sheetnames,
                    "sheets": sheets,
                    "formulas": formulas,
                }

            if ext == ".docx":
                d = docx.Document(file_path)
                all_lines = [p.text for p in d.paragraphs]
                paragraphs = [
                    {"index": i, "text": p.text, "style": getattr(p.style, "name", "Normal")}
                    for i, p in enumerate(d.paragraphs)
                    if p.text.strip()
                ]
                return {"format": "docx", "text": "\n".join(all_lines), "paragraphs": paragraphs}

            if ext == ".pptx":
                if not PPTX_AVAILABLE:
                    return {}
                prs = pptx.Presentation(file_path)
                slides: List[Dict[str, Any]] = []
                for idx, slide in enumerate(prs.slides):
                    title = ""
                    try:
                        if slide.shapes.title is not None:
                            title = slide.shapes.title.text or ""
                    except Exception:
                        pass
                    parts: List[str] = []
                    for ph in slide.placeholders:
                        try:
                            if ph.placeholder_format.idx == 0:
                                continue
                        except Exception:
                            pass
                        if ph.has_text_frame and ph.text_frame.text:
                            parts.append(ph.text_frame.text)
                    slides.append({
                        "slide_number": idx + 1,
                        "title": title,
                        "content": "\n".join(parts),
                    })
                return {"format": "pptx", "slides": slides}
        except Exception as e:
            logger.debug(f"Structured read skipped for {file_path}: {e}")
        return {}

    # ------------------------------------------------------------------
    # Canvas → File (user edits)
    # ------------------------------------------------------------------

    @staticmethod
    def _set_paragraph_text(paragraph, text: str) -> None:
        """Replace a paragraph's text while keeping its style and first-run
        formatting where possible."""
        if paragraph.runs:
            paragraph.runs[0].text = text
            for run in paragraph.runs[1:]:
                run.text = ""
        else:
            paragraph.add_run(text)

    def _sync_docx_in_place(self, file_path: str, content: str) -> None:
        """Line↔paragraph sync that PRESERVES document structure.

        The old implementation built a fresh Document() from plain text lines,
        destroying headings, tables, images and styles. Here we update the
        text of existing paragraphs in place (keeping their styles), blank any
        surplus paragraphs, and append only genuinely new trailing lines —
        tables and images between paragraphs survive every edit.
        """
        try:
            d = docx.Document(file_path)
        except Exception:
            # Target exists but isn't parseable as .docx (e.g. placeholder
            # bytes written by tooling). There is no structure to preserve —
            # fall back to a fresh document carrying the new lines.
            d = docx.Document()
        lines = content.split("\n")
        paras = list(d.paragraphs)

        for i, p in enumerate(paras):
            new_text = lines[i] if i < len(lines) else ""
            if p.text != new_text:
                self._set_paragraph_text(p, new_text)

        if len(lines) > len(paras):
            for line in lines[len(paras):]:
                d.add_paragraph(line)

        d.save(file_path)

    def sync_canvas_to_file(
        self,
        canvas_id: str,
        file_path: str,
        user_id: str,
        edit_type: str,  # 'cell' | 'document' | 'slide' | 'add_slide'
        data: Dict[str, Any]
    ) -> Dict[str, Any]:
        """Apply Canvas user edit back to filesystem document."""
        # R53: contain the path BEFORE any read/write — every OfficeService
        # entry point validates, but this service previously wrote docx files
        # via doc.save() on an unvalidated path (arbitrary file overwrite).
        try:
            file_path = _validate_office_path(file_path)
        except ValueError as e:
            return {"success": False, "error": str(e)}

        if not os.path.exists(file_path):
            return {"success": False, "error": f"File not found: {file_path}"}

        ext = Path(file_path).suffix.lower()
        try:
            if ext == ".xlsx" and edit_type == "cell":
                cell_path = data.get("cell_path")
                value = data.get("value")
                is_formula = data.get("is_formula", False)

                if not cell_path:
                    return {"success": False, "error": "cell_path required for Excel cell update"}

                res = self.office.excel.write_cell(
                    file_path=file_path,
                    cell_path=cell_path,
                    value=value,
                    is_formula=is_formula
                )
                if not res.get("success"):
                    return res

            elif ext == ".docx" and edit_type == "document":
                content = data.get("content")
                if content is None:
                    return {"success": False, "error": "content required for document update"}
                # In-place paragraph sync — preserves headings, tables, images
                # and paragraph styles (the old rebuild lost all of them).
                self._sync_docx_in_place(file_path, str(content))

            elif ext == ".pptx" and edit_type == "slide":
                res = self.office.pptx.modify_slides(
                    file_path=file_path,
                    action="update_slide",
                    options={
                        "slide_number": data.get("slide_number"),
                        "title": data.get("title"),
                        "content": data.get("content"),
                    },
                )
                if not res.get("success"):
                    return res

            elif ext == ".pptx" and edit_type == "add_slide":
                options = {
                    "title": data.get("title"),
                    "content": data.get("content"),
                }
                if data.get("layout_idx") is not None:
                    options["layout_idx"] = data.get("layout_idx")
                res = self.office.pptx.modify_slides(
                    file_path=file_path, action="add_slide", options=options,
                )
                if not res.get("success"):
                    return res

            else:
                return {"success": False, "error": f"Unsupported sync edit type: {edit_type} for extension {ext}"}

            # Snapshot AFTER the write so callers/UI get fresh content.
            content_snapshot = self._read_structured_content(file_path)
            component, canvas_type = OFFICE_COMPONENT_MAP.get(ext, ("office_word", "docs"))

            # Keep any persisted Canvas row's content fresh so reloading
            # /canvas/{id} (or a phone fetching GET /api/canvas/{id}) shows the
            # edited file — including the HTML preview, since mobile clients
            # have no canvas WebSocket subscription and load via REST only.
            # (No-op when the canvas row doesn't exist or the session isn't a
            # real DB session.) The broadcast below re-persists with the fresh
            # render; here we refresh the snapshot alone.
            self._refresh_canvas_row(canvas_id, content_snapshot, user_id)

            # After saving the file, push updated preview back to all subscribers
            self.broadcast_file_update(canvas_id, file_path, user_id)

            return {
                "success": True,
                "message": f"Successfully synchronized canvas {canvas_id} changes to {file_path}",
                "component": component,
                "canvas_type": canvas_type,
                "content": self._read_structured_content(file_path),
            }

        except Exception as e:
            logger.error(f"Failed to sync canvas to file: {e}")
            return {"success": False, "error": "Failed to sync canvas to file"}

    # ------------------------------------------------------------------
    # Canvas row persistence (present / refresh)
    # ------------------------------------------------------------------

    def ensure_canvas_for_file(
        self,
        canvas_id: Optional[str],
        file_path: str,
        user_id: str,
        title: Optional[str] = None,
    ) -> Dict[str, Any]:
        """Get-or-create the persistent Canvas row bound to an office file.

        Reuses the newest active canvas already bound to the same file when no
        explicit canvas_id is given (so repeated presents don't spawn duplicate
        canvases). Returns {id, created(bool), canvas_type, component, content}.
        Raises ValueError when the path escapes the office directory.
        """
        contained = _validate_office_path(file_path)
        ext = Path(contained).suffix.lower()
        component, canvas_type = OFFICE_COMPONENT_MAP.get(ext, ("office_word", "docs"))
        display_name = title or Path(contained).name

        row: Optional[Canvas] = None
        if canvas_id:
            found = self.db.query(Canvas).filter(Canvas.id == canvas_id).first()
            row = found if isinstance(found, Canvas) else None

        if row is None and not canvas_id:
            candidates = (
                self.db.query(Canvas)
                .filter(Canvas.status == "active")
                .order_by(Canvas.updated_at.desc())
                .limit(100)
                .all()
            )
            for cand in candidates:
                meta = cand.content if isinstance(cand.content, dict) else {}
                if meta.get("office_file") == contained:
                    row = cand
                    break

        created = False
        if row is None:
            row = Canvas(
                id=canvas_id or f"canvas_{uuid.uuid4().hex[:12]}",
                tenant_id="default",
                created_by=user_id or "system",
                name=display_name,
                canvas_type=canvas_type,
                status="active",
                content={"office_file": contained, "format": ext.lstrip(".")},
            )
            self.db.add(row)
            created = True

        snapshot = dict(self._read_structured_content(contained))
        meta = dict(row.content if isinstance(row.content, dict) else {})
        meta.update(snapshot)
        meta["office_file"] = contained
        meta["format"] = ext.lstrip(".")
        row.content = meta
        row.name = display_name
        row.last_edited_by = user_id or None
        row.last_edited_at = datetime.now(timezone.utc)
        self.db.commit()

        return {
            "id": row.id,
            "created": created,
            "canvas_type": canvas_type,
            "component": component,
            "content": meta,
        }

    def _refresh_canvas_row(
        self, canvas_id: str, content_snapshot: Dict[str, Any], user_id: str
    ) -> None:
        """Best-effort content refresh of an existing bound Canvas row."""
        try:
            found = self.db.query(Canvas).filter(Canvas.id == canvas_id).first()
            if not isinstance(found, Canvas):
                return
            meta = dict(found.content if isinstance(found.content, dict) else {})
            meta.update(content_snapshot)
            found.content = meta
            found.last_edited_by = user_id or None
            found.last_edited_at = datetime.now(timezone.utc)
            self.db.commit()
        except Exception as e:
            logger.debug(f"Canvas row refresh skipped for {canvas_id}: {e}")

    # ------------------------------------------------------------------
    # Agent-edit fan-out (file → canvases)
    # ------------------------------------------------------------------

    def notify_file_canvases(
        self, file_path: str, user_id: str, limit: int = 25
    ) -> List[str]:
        """Broadcast updated renders to every active canvas bound to a file.

        Called after agent-side writes (tools/office_tool.py) so user-facing
        canvases reflect agent edits. Returns the notified canvas ids.
        """
        try:
            contained = _validate_office_path(file_path)
        except ValueError:
            return []

        db = self.db
        owns_session = False
        if not isinstance(db, Session):
            from core.database import get_db_session as _get_session
            db = _get_session().__enter__()
            owns_session = True

        try:
            candidates = (
                db.query(Canvas)
                .filter(Canvas.status == "active")
                .order_by(Canvas.updated_at.desc())
                .limit(200)
                .all()
            )
            targets = [
                c for c in candidates
                if isinstance(c.content, dict) and c.content.get("office_file") == contained
            ][: max(0, limit)]
            for canvas in targets:
                self.broadcast_file_update(canvas.id, contained, user_id)
            return [c.id for c in targets]
        finally:
            if owns_session:
                try:
                    db.close()
                except Exception:
                    pass

    # ------------------------------------------------------------------
    # File → Canvas broadcast
    # ------------------------------------------------------------------

    def broadcast_file_update(self, canvas_id: str, file_path: str, user_id: str):
        """Broadcast updated document HTML render to the Canvas WebSocket subscribers."""
        # R53: contain the path — this reads the file (render + memory
        # ingestion) and stores its HTML in CanvasAudit/WS state the caller
        # controls; without containment any existing office file is readable.
        try:
            file_path = _validate_office_path(file_path)
        except ValueError:
            return

        try:
            render_res = self.office.renderer.render_to_html(file_path)
            # A failed HTML render (e.g. mammoth missing for docx) must NOT
            # abort the broadcast — the structured snapshot is independent of
            # the render, and bailing here left the canvas with no audit row
            # (making /api/canvas/{id} 404) and no WS update. Proceed with
            # html=None; the preview pane degrades gracefully.
            html = render_res.get("html") if render_res.get("success") else None

            ext = Path(file_path).suffix.lower()
            component, canvas_type = OFFICE_COMPONENT_MAP.get(ext, ("office_word", "docs"))
            component_type = "data_grid" if ext == ".xlsx" else "rich_editor"

            # Structured snapshot rides along so editable UIs can render real
            # inputs instead of scraping HTML.
            content_snapshot = self._read_structured_content(file_path)
            payload_data: Dict[str, Any] = dict(content_snapshot)
            payload_data.update({
                "html": html,
                "file_path": file_path,
                "format": ext.lstrip("."),
            })

            # Create updated audit record to update state history.
            # `content` nests the FULL client payload (structured snapshot +
            # html + file binding): read_canvas serves details_json.content
            # verbatim, so GET /api/canvas/{id} gives REST-only clients
            # (mobile has no canvas WS subscription) everything they need to
            # render an office canvas — including the /canvas/{id} web page
            # reload path.
            audit_content = dict(payload_data)
            audit_content["office_file"] = file_path
            audit_content["title"] = Path(file_path).name
            audit = CanvasAudit(
                id=str(uuid.uuid4()),
                tenant_id="default",
                user_id=user_id,
                canvas_id=canvas_id,
                action_type="update",
                canvas_type=canvas_type,
                details_json={
                    "canvas_type": canvas_type,
                    "component_type": component_type,
                    "component": component,
                    "file_path": file_path,
                    "html": html,
                    "title": Path(file_path).name,
                    "content": audit_content,
                }
            )
            self.db.add(audit)
            self.db.commit()

            # Persist the fresh render into the bound Canvas row so REST-only
            # clients (mobile has no canvas WS subscription) get a preview
            # from plain GET /api/canvas/{id}. Best-effort.
            self._refresh_canvas_row(canvas_id, {"html": html}, user_id)

            # Ingest the updated document content into Atom memory so the agent
            # remembers quotes/POs/price-lists/invoices it just generated or edited.
            # Fire-and-forget (async) to avoid blocking the Canvas update.
            import asyncio
            ingest_coro = self._ingest_document_to_memory(file_path, user_id)
            try:
                asyncio.create_task(ingest_coro)
            except RuntimeError:
                # No running loop (rare sync caller) — close the abandoned
                # coroutine (else it leaks with a "never awaited" warning) and
                # fall back to running it directly.
                ingest_coro.close()
                self._ingest_document_to_memory_sync(file_path, user_id)

            # Push live update via WebSocket manager. Deliver on BOTH channels:
            # canvas:{id} (page-scoped subscribers) AND user:{uid} (the working
            # channel chat panels actually listen on — office pushes were
            # dead-lettered before the user-channel leg existed).
            message = {
                "type": "canvas:update",
                "data": {
                    "action": "update",
                    "canvas_id": canvas_id,
                    "canvas_type": canvas_type,
                    "component": component,
                    "data": payload_data,
                }
            }
            for channel in (f"canvas:{canvas_id}", f"user:{user_id}"):
                try:
                    broadcast_coro = ws_manager.broadcast(channel, dict(message))
                    asyncio.create_task(broadcast_coro)
                except RuntimeError:
                    # No running loop — close the abandoned coroutine, skip async
                    # broadcast (rare sync caller).
                    try:
                        broadcast_coro.close()
                    except Exception:
                        pass

            # The channel-aware ws_manager only reaches clients connected via
            # the (currently unmounted) root websocket_routes /ws. Clients
            # actually connect through api/websocket_routes, which uses
            # notification_manager (workspace-wide, no channels) — mirror the
            # message there too. Clients filter by canvas_id, so a workspace-
            # wide fanout is harmless.
            try:
                from core.notification_manager import notification_manager
                nm_coro = notification_manager.broadcast(message, "default")
                asyncio.create_task(nm_coro)
            except Exception as e:
                try:
                    nm_coro.close()
                except Exception:
                    pass
                logger.debug(f"notification_manager mirror broadcast skipped: {e}")

        except Exception as e:
            logger.error(f"Failed to broadcast file update: {e}")

    async def _ingest_document_to_memory(self, file_path: str, user_id: str) -> bool:
        """Read a document's content and ingest it into Atom memory (async).

        Reuses AutoDocumentIngestionService.process_file_bytes so the same
        parse → redact → LanceDB + knowledge-graph path used for cloud-drive
        files applies to locally-edited Office documents. Failures are
        non-fatal (best-effort) and only logged.
        """
        try:
            content = self._read_file_bytes(file_path)
            if not content:
                return False
            from core.auto_document_ingestion import AutoDocumentIngestionService

            ingestor = AutoDocumentIngestionService()
            result = await ingestor.process_file_bytes(
                content=content,
                file_name=Path(file_path).name,
                source="office_canvas",
                user_id=user_id,
            )
            status = result.get("status")
            if status == "ingested":
                logger.info(
                    f"Office document ingested to memory: {Path(file_path).name} "
                    f"({result.get('chars_ingested', 0)} chars)"
                )
            return status in ("ingested", "skipped")
        except Exception as e:
            logger.debug(f"Office→memory ingestion skipped for {file_path}: {e}")
            return False

    def _ingest_document_to_memory_sync(self, file_path: str, user_id: str) -> bool:
        """Synchronous fallback for Office→memory ingestion (no running event loop)."""
        try:
            import asyncio
            content = self._read_file_bytes(file_path)
            if not content:
                return False
            from core.auto_document_ingestion import AutoDocumentIngestionService

            ingestor = AutoDocumentIngestionService()
            loop = asyncio.new_event_loop()
            try:
                result = loop.run_until_complete(
                    ingestor.process_file_bytes(
                        content=content,
                        file_name=Path(file_path).name,
                        source="office_canvas",
                        user_id=user_id,
                    )
                )
            finally:
                # Always close the throwaway loop — leaking it keeps an open
                # selector FD per sync-context office edit.
                loop.close()
            return result.get("status") in ("ingested", "skipped")
        except Exception as e:
            logger.debug(f"Office→memory sync ingestion skipped for {file_path}: {e}")
            return False

    @staticmethod
    def _read_file_bytes(file_path: str) -> Optional[bytes]:
        """Read file bytes if the file exists and is non-empty."""
        try:
            if not os.path.exists(file_path):
                return None
            with open(file_path, "rb") as f:
                content = f.read()
            return content or None
        except Exception:
            return None
