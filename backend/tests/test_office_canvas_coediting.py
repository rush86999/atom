"""
Office Canvas Co-Editing — user edits on the canvas must modify the real
file, and agent edits to the file must reach every open canvas (Red-Green-
Refactor companion to the #39 co-editing loop).

Gaps covered here:
  1. docx canvas→file sync was a destructive plain-text rebuild (fresh
     Document()) — styles, tables and images were lost. Must now be an
     in-place paragraph sync that preserves document structure.
  2. pptx had NO sync path at all ("Unsupported sync edit type"). Slide-text
     edits and slide creation must round-trip to the file.
  3. File→canvas broadcast carried only rendered HTML with a generic
     "office_preview" component on a dead-lettered channel. It must carry a
     STRUCTURED content snapshot (grid cells / paragraphs / slides), a
     renderable component name (office_excel/office_word/office_pptx), and
     fan out to BOTH canvas:{id} and user:{uid} channels.
  4. /present never persisted a Canvas row — the presented office canvas was
     unreachable at /canvas/{id} after refresh (and CanvasAudit rows dangled
     against a missing parent on FK-enforcing databases). Present must
     create/reuse a DB Canvas row bound to the file.
  5. Agent writes through tools/office_tool.py silently skipped any open
     canvas — writes must notify canvases bound to the file.
"""

import asyncio
import uuid
from unittest.mock import AsyncMock, MagicMock, patch

import openpyxl
import pytest
from fastapi import FastAPI
from fastapi.testclient import TestClient
from sqlalchemy import create_engine
from sqlalchemy.orm import sessionmaker
from sqlalchemy.pool import StaticPool

from core.auth import get_current_user as auth_get_current_user
from core.database import Base, get_db_session


# ----------------------------------------------------------------------------
# Fixtures
# ----------------------------------------------------------------------------

@pytest.fixture()
def office_root(tmp_path, monkeypatch):
    office_dir = tmp_path / "office"
    office_dir.mkdir()
    monkeypatch.setenv("ATOM_OFFICE_DIR", str(office_dir))
    return office_dir


@pytest.fixture()
def test_db():
    engine = create_engine(
        "sqlite://",
        connect_args={"check_same_thread": False},
        poolclass=StaticPool,
    )
    Base.metadata.create_all(engine)
    factory = sessionmaker(bind=engine, autoflush=False, autocommit=False)
    session = factory()
    try:
        yield session
    finally:
        session.close()


def _make_xlsx(path):
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws["A1"] = "Item"
    ws["B1"] = "Qty"
    ws["A2"] = "Widget"
    ws["B2"] = 4
    wb.save(path)


@pytest.fixture()
def xlsx_file(office_root):
    p = office_root / "book.xlsx"
    _make_xlsx(p)
    return p


def _make_docx(path):
    import docx

    doc = docx.Document()
    doc.add_paragraph("Quarterly Report", style="Heading 1")
    doc.add_paragraph("Intro paragraph.")
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "K"
    table.cell(0, 1).text = "V"
    table.cell(1, 0).text = "Revenue"
    table.cell(1, 1).text = "$100"
    doc.add_paragraph("Closing note.", style="Intense Quote")
    doc.save(path)


@pytest.fixture()
def docx_file(office_root):
    p = office_root / "report.docx"
    _make_docx(p)
    return p


def _make_pptx(path):
    import pptx

    prs = pptx.Presentation()
    layout = prs.slide_layouts[1]  # Title + Content
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Original Title"
    slide.placeholders[1].text = "Original Body"
    prs.save(path)


@pytest.fixture()
def pptx_file(office_root):
    p = office_root / "deck.pptx"
    _make_pptx(p)
    return p


# ----------------------------------------------------------------------------
# 1+6. Structured content snapshots & sync results
# ----------------------------------------------------------------------------

class TestStructuredContent:
    def test_xlsx_snapshot_has_grid(self, xlsx_file):
        from core.office_sync_service import OfficeSyncService

        snap = OfficeSyncService._read_structured_content(str(xlsx_file))
        assert snap["format"] == "xlsx"
        assert snap["active_sheet"] == "Sheet1"
        sheet = snap["sheets"][0]
        assert sheet["name"] == "Sheet1"
        assert ["Item", "Qty"] in [list(map(str, r[:2])) for r in sheet["rows"]]

    def test_docx_snapshot_joins_all_paragraphs(self, docx_file):
        from core.office_sync_service import OfficeSyncService

        snap = OfficeSyncService._read_structured_content(str(docx_file))
        assert snap["format"] == "docx"
        assert "Quarterly Report" in snap["text"]
        assert "Closing note." in snap["text"]
        styles = [p["style"] for p in snap["paragraphs"]]
        assert "Heading 1" in styles

    def test_pptx_snapshot_lists_slides(self, pptx_file):
        from core.office_sync_service import OfficeSyncService

        snap = OfficeSyncService._read_structured_content(str(pptx_file))
        assert snap["format"] == "pptx"
        assert snap["slides"][0]["slide_number"] == 1
        assert snap["slides"][0]["title"] == "Original Title"
        assert "Original Body" in snap["slides"][0]["content"]

    def test_sync_result_includes_snapshot(self, xlsx_file):
        """A successful canvas→file sync must return the fresh content so the
        UI can update its own state immediately."""
        from core.office_sync_service import OfficeSyncService

        svc = OfficeSyncService(MagicMock())
        res = svc.sync_canvas_to_file(
            canvas_id="c1",
            file_path=str(xlsx_file),
            user_id="u1",
            edit_type="cell",
            data={"cell_path": "/Sheet1/B2", "value": 10},
        )
        assert res["success"] is True
        assert res["content"]["format"] == "xlsx"
        assert res["component"] == "office_excel"
        flat = [str(v) for row in res["content"]["sheets"][0]["rows"] for v in row]
        assert "10" in flat


# ----------------------------------------------------------------------------
# 2. Non-destructive docx sync
# ----------------------------------------------------------------------------

class TestDocxNonDestructiveSync:
    def test_edit_preserves_table_and_styles(self, docx_file):
        import docx as docx_lib

        from core.office_sync_service import OfficeSyncService

        res = OfficeSyncService(MagicMock()).sync_canvas_to_file(
            canvas_id="c1",
            file_path=str(docx_file),
            user_id="u1",
            edit_type="document",
            data={"content": "Quarterly Report\nEDITED INTRO.\n\nClosing note."},
        )
        assert res["success"] is True

        doc = docx_lib.Document(str(docx_file))
        texts = [p.text for p in doc.paragraphs]
        assert "EDITED INTRO." in texts
        # Heading style survived (old code rebuilt a bare Document()).
        assert doc.paragraphs[0].style.name == "Heading 1"
        # Table survived the edit.
        assert len(doc.tables) == 1
        assert doc.tables[0].cell(1, 1).text == "$100"

    def test_appended_lines_added_with_paragraph(self, docx_file):
        import docx as docx_lib

        from core.office_sync_service import OfficeSyncService

        res = OfficeSyncService(MagicMock()).sync_canvas_to_file(
            canvas_id="c1",
            file_path=str(docx_file),
            user_id="u1",
            edit_type="document",
            data={
                "content": (
                    "Quarterly Report\nIntro paragraph.\nBRAND NEW LINE"
                )
            },
        )
        assert res["success"] is True
        doc = docx_lib.Document(str(docx_file))
        assert "BRAND NEW LINE" in [p.text for p in doc.paragraphs]


# ----------------------------------------------------------------------------
# 3. PPTX sync path (previously unsupported)
# ----------------------------------------------------------------------------

class TestPptxSync:
    def test_slide_text_edit_round_trips(self, pptx_file):
        import pptx as pptx_lib

        from core.office_sync_service import OfficeSyncService

        res = OfficeSyncService(MagicMock()).sync_canvas_to_file(
            canvas_id="c1",
            file_path=str(pptx_file),
            user_id="u1",
            edit_type="slide",
            data={"slide_number": 1, "title": "New Title", "content": "New Body"},
        )
        assert res["success"] is True
        assert res["content"]["format"] == "pptx"

        prs = pptx_lib.Presentation(str(pptx_file))
        assert prs.slides[0].shapes.title.text == "New Title"
        assert prs.slides[0].placeholders[1].text == "New Body"

    def test_add_slide_creates_slide(self, pptx_file):
        import pptx as pptx_lib

        from core.office_sync_service import OfficeSyncService

        res = OfficeSyncService(MagicMock()).sync_canvas_to_file(
            canvas_id="c1",
            file_path=str(pptx_file),
            user_id="u1",
            edit_type="add_slide",
            data={"title": "Second", "content": "More"},
        )
        assert res["success"] is True
        prs = pptx_lib.Presentation(str(pptx_file))
        assert len(prs.slides) == 2
        assert prs.slides[1].shapes.title.text == "Second"


# ----------------------------------------------------------------------------
# 4. Broadcast: structured payload + dual-channel delivery
# ----------------------------------------------------------------------------

class TestBroadcastPayload:
    @pytest.mark.parametrize(
        "fixture_name,component,fmt",
        [
            ("xlsx_file", "office_excel", "xlsx"),
            ("docx_file", "office_word", "docx"),
            ("pptx_file", "office_pptx", "pptx"),
        ],
    )
    def test_component_names_and_structured_data(self, request, fixture_name, component, fmt):
        from core.office_sync_service import OfficeSyncService

        path = request.getfixturevalue(fixture_name)
        svc = OfficeSyncService(MagicMock())
        with patch.object(svc.office.renderer, "render_to_html",
                          return_value={"success": True, "html": "<b>x</b>"}), \
             patch.object(svc, "_ingest_document_to_memory",
                          new=AsyncMock(return_value=True)), \
             patch("core.office_sync_service.ws_manager.broadcast",
                   new=AsyncMock()) as wsb:
            svc.broadcast_file_update("c1", str(path), "u1")

        # Fire-and-forget broadcast: the WS manager records each publish at
        # coroutine-creation time (a running loop delivers it in prod).
        assert wsb.call_count >= 1
        payloads = [c.args[1] for c in wsb.call_args_list]
        target = next(p for p in payloads if p["data"].get("component") == component)
        inner = target["data"]["data"]
        assert inner["format"] == fmt
        assert inner["html"] == "<b>x</b>"
        assert inner["file_path"]
        # Structured snapshot rides along for editable UIs.
        key = {"xlsx": "sheets", "docx": "text", "pptx": "slides"}[fmt]
        assert key in inner

    def test_broadcast_reaches_both_channels(self, xlsx_file):
        from core.office_sync_service import OfficeSyncService

        svc = OfficeSyncService(MagicMock())
        with patch.object(svc.office.renderer, "render_to_html",
                          return_value={"success": True, "html": "h"}), \
             patch.object(svc, "_ingest_document_to_memory",
                          new=AsyncMock(return_value=True)), \
             patch("core.office_sync_service.ws_manager.broadcast",
                   new=AsyncMock()) as wsb:
            svc.broadcast_file_update("c1", str(xlsx_file), "u1")

        channels = [c.args[0] for c in wsb.call_args_list]
        assert "canvas:c1" in channels, "canvas-scoped subscribers got nothing"
        assert "user:u1" in channels, (
            "office updates were dead-lettered: no client ever subscribes "
            "to canvas:* unless it can also arrive on user:{uid}"
        )


# ----------------------------------------------------------------------------
# 5. Canvas persistence on present + reuse
# ----------------------------------------------------------------------------

class TestPresentPersistence:
    def _client(self, test_db):
        from api.office_routes import router, get_db_session_dep

        app = FastAPI()
        app.include_router(router)
        app.dependency_overrides[auth_get_current_user] = lambda: MagicMock(
            id="u-pres", email="p@example.com"
        )
        app.dependency_overrides[get_db_session_dep] = lambda: test_db
        return TestClient(app, raise_server_exceptions=False)

    def test_present_creates_bound_canvas_row(self, test_db, xlsx_file):
        from core.models import Canvas

        resp = self._client(test_db).post(
            "/present",
            json={"file_path": str(xlsx_file), "user_id": "ignored"},
        )
        assert resp.status_code == 200
        canvas_id = resp.json()["canvas_id"]

        row = test_db.query(Canvas).filter(Canvas.id == canvas_id).first()
        assert row is not None, "/present did not persist a Canvas row"
        assert row.content.get("office_file") == str(xlsx_file)
        assert row.canvas_type == "sheets"

    def test_present_reuses_existing_binding(self, test_db, xlsx_file):
        client = self._client(test_db)
        first = client.post("/present", json={"file_path": str(xlsx_file), "user_id": "x"}).json()
        second = client.post("/present", json={"file_path": str(xlsx_file), "user_id": "x"}).json()
        assert first["canvas_id"] == second["canvas_id"], (
            "every present minted a new canvas — office files would spawn "
            "unbounded duplicate canvases"
        )

    def test_explicit_canvas_id_respected(self, test_db, xlsx_file):
        from core.models import Canvas

        cid = f"canvas_{uuid.uuid4().hex[:12]}"
        resp = self._client(test_db).post(
            "/present",
            json={"file_path": str(xlsx_file), "user_id": "x", "canvas_id": cid},
        )
        assert resp.json()["canvas_id"] == cid
        assert test_db.query(Canvas).filter(Canvas.id == cid).first() is not None

    def test_canvas_types_per_format(self, test_db, office_root):
        from core.models import Canvas

        docx_p = office_root / "d.docx"
        pptx_p = office_root / "d.pptx"
        _make_docx(docx_p)
        _make_pptx(pptx_p)

        client = self._client(test_db)
        r1 = client.post("/present", json={"file_path": str(docx_p), "user_id": "x"}).json()
        r2 = client.post("/present", json={"file_path": str(pptx_p), "user_id": "x"}).json()

        t1 = test_db.query(Canvas).filter(Canvas.id == r1["canvas_id"]).first().canvas_type
        t2 = test_db.query(Canvas).filter(Canvas.id == r2["canvas_id"]).first().canvas_type
        assert t1 == "docs"
        assert t2 == "presentation"


# ----------------------------------------------------------------------------
# 6. notify_file_canvases — agent edits reach open canvases
# ----------------------------------------------------------------------------

class TestNotifyFileCanvases:
    def test_agent_write_notifies_bound_canvas(self, test_db, xlsx_file):
        from core.models import Canvas
        from core.office_sync_service import OfficeSyncService

        cid = f"canvas_{uuid.uuid4().hex[:12]}"
        test_db.add(Canvas(
            id=cid,
            tenant_id="default",
            created_by=None or "u1",
            name="book.xlsx",
            canvas_type="sheets",
            content={"office_file": str(xlsx_file), "format": "xlsx"},
        ))
        test_db.commit()

        svc = OfficeSyncService(test_db)
        notified = None
        with patch.object(svc.office.renderer, "render_to_html",
                          return_value={"success": True, "html": "h"}), \
             patch.object(svc, "_ingest_document_to_memory",
                          new=AsyncMock(return_value=True)), \
             patch("core.office_sync_service.ws_manager.broadcast",
                   new=AsyncMock()):
            notified = svc.notify_file_canvases(str(xlsx_file), "agent-u")

        assert cid in notified

    def test_unbound_files_notify_nothing(self, test_db, tmp_path):
        from core.office_sync_service import OfficeSyncService

        stray = tmp_path / "elsewhere.xlsx"
        _make_xlsx(stray)
        svc = OfficeSyncService(test_db)
        assert svc.notify_file_canvases(str(stray), "u1") == []


# ----------------------------------------------------------------------------
# 7. Agent tool writes schedule canvas notification
# ----------------------------------------------------------------------------

class TestAgentToolNotifyHook:
    def _run_with_flush(self, coro_factory):
        async def main():
            result = await coro_factory()
            await asyncio.sleep(0)
            await asyncio.sleep(0)
            return result
        return asyncio.run(main())

    def test_write_excel_cell_notifies(self, monkeypatch, xlsx_file):
        import tools.office_tool as ot

        calls = []

        async def fake_notify(fp, uid):
            calls.append((fp, uid))

        monkeypatch.setattr(ot, "_notify_canvases", fake_notify)

        def factory():
            return ot.write_excel_cell(
                user_id="u1", file_path=str(xlsx_file),
                cell_path="/Sheet1/A2", value="Updated",
            )

        self._run_with_flush(factory)
        assert calls and calls[0][0] == str(xlsx_file)

    def test_failed_write_does_not_notify(self, monkeypatch, tmp_path):
        import tools.office_tool as ot

        calls = []

        async def fake_notify(fp, uid):
            calls.append(fp)

        monkeypatch.setattr(ot, "_notify_canvases", fake_notify)
        outside = tmp_path / "outside" / "book.xlsx"
        outside.parent.mkdir()

        def factory():
            return ot.write_excel_cell(
                user_id="u1",
                file_path=str(outside),  # out of scope → access denied
                cell_path="/Sheet1/A1", value=1,
            )

        self._run_with_flush(factory)
        assert calls == []

    def test_modify_word_notifies(self, monkeypatch, docx_file):
        import tools.office_tool as ot

        calls = []

        async def fake_notify(fp, uid):
            calls.append(fp)

        monkeypatch.setattr(ot, "_notify_canvases", fake_notify)

        def factory():
            return ot.modify_word_document(
                user_id="u1", file_path=str(docx_file),
                action="append", content="Appended line",
            )

        self._run_with_flush(factory)
        assert calls == [str(docx_file)]

    def test_modify_pptx_notifies(self, monkeypatch, pptx_file):
        import tools.office_tool as ot

        calls = []

        async def fake_notify(fp, uid):
            calls.append(fp)

        monkeypatch.setattr(ot, "_notify_canvases", fake_notify)

        def factory():
            return ot.modify_pptx_slides(
                user_id="u1", file_path=str(pptx_file),
                action="add_slide", title="T", content="C",
            )

        self._run_with_flush(factory)
        assert calls == [str(pptx_file)]


# ----------------------------------------------------------------------------
# 8. sync-update keeps an existing Canvas row's content fresh
# ----------------------------------------------------------------------------

class TestSyncUpdatesCanvasRow:
    def test_cell_edit_refreshes_canvas_content(self, test_db, xlsx_file):
        from core.models import Canvas
        from core.office_sync_service import OfficeSyncService

        cid = f"canvas_{uuid.uuid4().hex[:12]}"
        test_db.add(Canvas(
            id=cid, tenant_id="default", created_by="u1", name="book.xlsx",
            canvas_type="sheets",
            content={"office_file": str(xlsx_file), "format": "xlsx"},
        ))
        test_db.commit()

        OfficeSyncService(test_db).sync_canvas_to_file(
            canvas_id=cid, file_path=str(xlsx_file), user_id="u1",
            edit_type="cell", data={"cell_path": "/Sheet1/B2", "value": 42},
        )

        row = test_db.query(Canvas).filter(Canvas.id == cid).first()
        flat = [str(v) for r in row.content["sheets"][0]["rows"] for v in r]
        assert "42" in flat, "open canvas kept stale content after a file edit"

    def test_broadcast_persists_html_for_rest_clients(self, test_db, xlsx_file):
        """Mobile has no canvas WebSocket subscription — it loads canvases via
        plain GET /api/canvas/{id}, which serves the latest CanvasAudit row's
        details_json.content. That dict must carry the HTML preview, the
        structured snapshot and the office_file binding."""
        from core.models import Canvas, CanvasAudit
        from core.office_sync_service import OfficeSyncService

        cid = f"canvas_{uuid.uuid4().hex[:12]}"
        test_db.add(Canvas(
            id=cid, tenant_id="default", created_by="u1", name="book.xlsx",
            canvas_type="sheets",
            content={"office_file": str(xlsx_file), "format": "xlsx"},
        ))
        test_db.commit()

        svc = OfficeSyncService(test_db)
        with patch.object(svc.office.renderer, "render_to_html",
                          return_value={"success": True, "html": "<table>preview</table>"}), \
             patch.object(svc, "_ingest_document_to_memory",
                          new=AsyncMock(return_value=True)), \
             patch("core.office_sync_service.ws_manager.broadcast", new=AsyncMock()):
            svc.broadcast_file_update(cid, str(xlsx_file), "u1")

        # Canvas row (web reload path)
        row = test_db.query(Canvas).filter(Canvas.id == cid).first()
        assert row.content.get("html") == "<table>preview</table>", (
            "REST-only clients got no preview: broadcast did not persist html"
        )

        # Audit trail (GET /api/canvas/{id} source of truth)
        audit = (test_db.query(CanvasAudit)
                 .filter(CanvasAudit.canvas_id == cid)
                 .order_by(CanvasAudit.created_at.desc())
                 .first())
        served = audit.details_json["content"]
        assert served["html"] == "<table>preview</table>"
        assert served["format"] == "xlsx"
        assert served["office_file"] == str(xlsx_file)
        assert "sheets" in served

    def test_magicmock_db_is_tolerated(self, xlsx_file):
        """Legacy callers pass MagicMock sessions — canvas refresh must be a
        no-op there, never a crash."""
        from core.office_sync_service import OfficeSyncService

        res = OfficeSyncService(MagicMock()).sync_canvas_to_file(
            canvas_id="c1", file_path=str(xlsx_file), user_id="u1",
            edit_type="cell", data={"cell_path": "/Sheet1/B2", "value": 7},
        )
        assert res["success"] is True
