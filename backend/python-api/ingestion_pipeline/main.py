from fastapi import FastAPI, HTTPException, BackgroundTasks
from pydantic import BaseModel
import uvicorn
import os
import logging
import datetime
from typing import Optional, List, Dict, Any
import asyncio

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

app = FastAPI(
    title="Ingestion Pipeline API",
    description="API for ingesting and processing various data sources",
    version="1.0.0"
)

class HealthResponse(BaseModel):
    status: str
    version: str
    timestamp: str
    processed_items: int = 0

class IngestionRequest(BaseModel):
    source_type: str  # "notion", "dropbox", "gdrive", "local", etc.
    source_config: Dict[str, Any]
    processing_mode: str = "incremental"  # "full" or "incremental"
    user_id: Optional[str] = None

class IngestionStatus(BaseModel):
    status: str  # "processing", "completed", "failed"
    total_items: int
    processed_items: int
    failed_items: int
    start_time: str
    end_time: Optional[str] = None
    error_message: Optional[str] = None

# In-memory storage for demonstration (replace with database in production)
ingestion_statuses = {}
processed_items_count = 0

@app.get("/")
async def root():
    return {"message": "Ingestion Pipeline API", "version": "1.0.0"}

@app.get("/healthz", response_model=HealthResponse)
async def health_check():
    return HealthResponse(
        status="ok",
        version="1.0.0",
        timestamp=datetime.datetime.now().isoformat(),
        processed_items=processed_items_count
    )

@app.get("/ingestion-status")
async def get_ingestion_status():
    """Get overall ingestion pipeline status"""
    return {
        "status": "operational",
        "total_processed": processed_items_count,
        "active_jobs": len([s for s in ingestion_statuses.values() if s["status"] == "processing"])
    }

@app.get("/ingestion-status/{job_id}", response_model=IngestionStatus)
async def get_job_status(job_id: str):
    """Get status of a specific ingestion job"""
    if job_id not in ingestion_statuses:
        raise HTTPException(status_code=404, detail="Job not found")

    status = ingestion_statuses[job_id]
    return IngestionStatus(**status)

@app.post("/start-ingestion")
async def start_ingestion(request: IngestionRequest, background_tasks: BackgroundTasks):
    """Start a new ingestion job"""
    job_id = f"ingestion_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}_{os.urandom(4).hex()}"

    # Initialize job status
    ingestion_statuses[job_id] = {
        "status": "processing",
        "total_items": 0,
        "processed_items": 0,
        "failed_items": 0,
        "start_time": datetime.datetime.now().isoformat(),
        "end_time": None,
        "error_message": None
    }

    # Start background processing
    background_tasks.add_task(process_ingestion, job_id, request)

    return {
        "job_id": job_id,
        "status": "started",
        "message": f"Ingestion job {job_id} started successfully"
    }

async def process_ingestion(job_id: str, request: IngestionRequest):
    """Background task to process ingestion"""
    try:
        logger.info(f"Starting ingestion job {job_id} for source: {request.source_type}")

        # Simulate processing different source types
        if request.source_type == "notion":
            items = await process_notion_source(request.source_config)
        elif request.source_type == "dropbox":
            items = await process_dropbox_source(request.source_config)
        elif request.source_type == "gdrive":
            items = await process_gdrive_source(request.source_config)
        elif request.source_type == "local":
            items = await process_local_source(request.source_config)
        else:
            raise ValueError(f"Unsupported source type: {request.source_type}")

        # Update job status
        ingestion_statuses[job_id]["total_items"] = len(items)

        # Process each item
        for i, item in enumerate(items):
            try:
                await process_item(item, request.user_id)
                ingestion_statuses[job_id]["processed_items"] += 1
                global processed_items_count
                processed_items_count += 1

                # Simulate processing time
                await asyncio.sleep(0.1)

            except Exception as e:
                logger.error(f"Error processing item {i}: {e}")
                ingestion_statuses[job_id]["failed_items"] += 1

        # Mark job as completed
        ingestion_statuses[job_id]["status"] = "completed"
        ingestion_statuses[job_id]["end_time"] = datetime.datetime.now().isoformat()

        logger.info(f"Completed ingestion job {job_id}. Processed: {ingestion_statuses[job_id]['processed_items']}, Failed: {ingestion_statuses[job_id]['failed_items']}")

    except Exception as e:
        logger.error(f"Error in ingestion job {job_id}: {e}")
        ingestion_statuses[job_id]["status"] = "failed"
        ingestion_statuses[job_id]["error_message"] = str(e)
        ingestion_statuses[job_id]["end_time"] = datetime.datetime.now().isoformat()

async def process_notion_source(config: Dict[str, Any]) -> List[Dict[str, Any]]:
    """Process Notion source using Notion API"""
    logger.info(f"Processing Notion source with config: {config}")

    try:
        # Extract configuration
        api_key = config.get("api_key")
        database_id = config.get("database_id")
        page_size = config.get("page_size", 100)

        if not api_key or not database_id:
            raise ValueError("Notion API key and database ID are required")

        # Query Notion database
        headers = {
            "Authorization": f"Bearer {api_key}",
            "Notion-Version": "2022-06-28",
            "Content-Type": "application/json"
        }

        url = f"https://api.notion.com/v1/databases/{database_id}/query"
        payload = {
            "page_size": page_size,
            "sorts": [{"timestamp": "last_edited_time", "direction": "descending"}]
        }

        async with aiohttp.ClientSession() as session:
            async with session.post(url, json=payload, headers=headers) as response:
                if response.status != 200:
                    error_text = await response.text()
                    raise Exception(f"Notion API error: {response.status} - {error_text}")

                data = await response.json()
                results = data.get("results", [])

                # Extract content from Notion pages
                items = []
                for page in results:
                    page_id = page.get("id", "")
                    properties = page.get("properties", {})

                    # Extract title from properties (adjust based on your database structure)
                    title = ""
                    for prop_name, prop_value in properties.items():
                        if prop_value.get("type") == "title" and prop_value.get("title"):
                            title = prop_value["title"][0].get("plain_text", "") if prop_value["title"] else ""
                            break

                    # Extract last edited time
                    last_edited = page.get("last_edited_time", "")

                    items.append({
                        "id": f"notion_{page_id}",
                        "content": title,
                        "metadata": {
                            "source": "notion",
                            "page_id": page_id,
                            "last_edited": last_edited,
                            "url": page.get("url", "")
                        }
                    })

                logger.info(f"Processed {len(items)} items from Notion")
                return items

    except Exception as e:
        logger.error(f"Error processing Notion source: {str(e)}")
        raise

async def process_dropbox_source(config: Dict[str, Any]) -> List[Dict[str, Any]]:
    """Process Dropbox source using Dropbox API"""
    logger.info(f"Processing Dropbox source with config: {config}")

    try:
        # Extract configuration
        access_token = config.get("access_token")
        folder_path = config.get("folder_path", "")
        file_types = config.get("file_types", [".txt", ".pdf", ".docx", ".md"])

        if not access_token:
            raise ValueError("Dropbox access token is required")

        # List files in Dropbox folder
        headers = {
            "Authorization": f"Bearer {access_token}",
            "Content-Type": "application/json"
        }

        url = "https://api.dropboxapi.com/2/files/list_folder"
        payload = {
            "path": folder_path,
            "recursive": True,
            "include_media_info": False,
            "include_deleted": False,
            "include_has_explicit_shared_members": False
        }

        async with aiohttp.ClientSession() as session:
            async with session.post(url, json=payload, headers=headers) as response:
                if response.status != 200:
                    error_text = await response.text()
                    raise Exception(f"Dropbox API error: {response.status} - {error_text}")

                data = await response.json()
                entries = data.get("entries", [])

                # Filter files by type and process them
                items = []
                for entry in entries:
                    if entry.get(".tag") == "file" and any(entry.get("name", "").endswith(ext) for ext in file_types):
                        file_path = entry.get("path_display", "")
                        file_id = entry.get("id", "")
                        file_name = entry.get("name", "")
                        modified = entry.get("client_modified", "")

                        # Download file content
                        download_url = "https://content.dropboxapi.com/2/files/download"
                        download_headers = {
                            "Authorization": f"Bearer {access_token}",
                            "Dropbox-API-Arg": json.dumps({"path": file_path})
                        }

                        async with session.post(download_url, headers=download_headers) as download_response:
                            if download_response.status == 200:
                                content = await download_response.text()

                                items.append({
                                    "id": f"dropbox_{file_id}",
                                    "content": content,
                                    "metadata": {
                                        "source": "dropbox",
                                        "file_path": file_path,
                                        "file_name": file_name,
                                        "modified": modified,
                                        "size": entry.get("size", 0)
                                    }
                                })

                logger.info(f"Processed {len(items)} files from Dropbox")
                return items

    except Exception as e:
        logger.error(f"Error processing Dropbox source: {str(e)}")
        raise

async def process_gdrive_source(config: Dict[str, Any]) -> List[Dict[str, Any]]:
    """Process Google Drive source using Google Drive API"""
    logger.info(f"Processing Google Drive source with config: {config}")

    try:
        # Extract configuration
        access_token = config.get("access_token")
        folder_id = config.get("folder_id", "root")
        mime_types = config.get("mime_types", ["text/plain", "application/pdf", "application/vnd.google-apps.document"])

        if not access_token:
            raise ValueError("Google Drive access token is required")

        # List files in Google Drive folder
        headers = {
            "Authorization": f"Bearer {access_token}",
            "Content-Type": "application/json"
        }

        # Build query for specific folder and file types
        query_parts = [f"'{folder_id}' in parents", "trashed = false"]
        if mime_types:
            mime_query = " or ".join([f"mimeType = '{mime}'" for mime in mime_types])
            query_parts.append(f"({mime_query})")

        query = " and ".join(query_parts)
        url = f"https://www.googleapis.com/drive/v3/files?q={urllib.parse.quote(query)}&fields=files(id,name,mimeType,modifiedTime,size)"

        async with aiohttp.ClientSession() as session:
            async with session.get(url, headers=headers) as response:
                if response.status != 200:
                    error_text = await response.text()
                    raise Exception(f"Google Drive API error: {response.status} - {error_text}")

                data = await response.json()
                files = data.get("files", [])

                # Process each file
                items = []
                for file in files:
                    file_id = file.get("id", "")
                    file_name = file.get("name", "")
                    mime_type = file.get("mimeType", "")
                    modified = file.get("modifiedTime", "")

                    # Download file content based on mime type
                    if mime_type == "application/vnd.google-apps.document":
                        # Export Google Docs as text
                        export_url = f"https://www.googleapis.com/drive/v3/files/{file_id}/export?mimeType=text/plain"
                        async with session.get(export_url, headers=headers) as export_response:
                            if export_response.status == 200:
                                content = await export_response.text()
                    else:
                        # Download regular files
                        download_url = f"https://www.googleapis.com/drive/v3/files/{file_id}?alt=media"
                        async with session.get(download_url, headers=headers) as download_response:
                            if download_response.status == 200:
                                content = await download_response.text()

                    items.append({
                        "id": f"gdrive_{file_id}",
                        "content": content,
                        "metadata": {
                            "source": "google_drive",
                            "file_id": file_id,
                            "file_name": file_name,
                            "mime_type": mime_type,
                            "modified": modified,
                            "size": file.get("size", 0)
                        }
                    })

                logger.info(f"Processed {len(items)} files from Google Drive")
                return items

    except Exception as e:
        logger.error(f"Error processing Google Drive source: {str(e)}")
        raise

async def process_local_source(config: Dict[str, Any]) -> List[Dict[str, Any]]:
    """Process local file system source"""
    logger.info(f"Processing local source with config: {config}")

    try:
        # Extract configuration
        directory_path = config.get("directory_path", ".")
        file_extensions = config.get("file_extensions", [".txt", ".md", ".pdf", ".docx", ".pptx", ".xlsx"])
        recursive = config.get("recursive", True)

        if not os.path.exists(directory_path):
            raise ValueError(f"Directory does not exist: {directory_path}")

        if not os.path.isdir(directory_path):
            raise ValueError(f"Path is not a directory: {directory_path}")

        # Scan directory for files
        items = []
        if recursive:
            walk_generator = os.walk(directory_path)
        else:
            walk_generator = [(directory_path, [], [f for f in os.listdir(directory_path) if os.path.isfile(os.path.join(directory_path, f))])]

        for root, dirs, files in walk_generator:
            for file_name in files:
                if any(file_name.lower().endswith(ext) for ext in file_extensions):
                    file_path = os.path.join(root, file_name)
                    relative_path = os.path.relpath(file_path, directory_path)

                    try:
                        # Read file content based on file type
                        if file_name.lower().endswith(('.txt', '.md', '.csv', '.json', '.xml', '.html', '.htm')):
                            # Text files
                            async with aiofiles.open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                                content = await f.read()
                        elif file_name.lower().endswith('.pdf'):
                            # PDF files
                            import pdfplumber
                            content = ""
                            with pdfplumber.open(file_path) as pdf:
                                for page in pdf.pages:
                                    page_text = page.extract_text()
                                    if page_text:
                                        content += page_text + "\n"
                        elif file_name.lower().endswith(('.docx', '.pptx', '.xlsx')):
                            # Office documents
                            if file_name.lower().endswith('.docx'):
                                from docx import Document
                                doc = Document(file_path)
                                content = "\n".join([paragraph.text for paragraph in doc.paragraphs])
                            elif file_name.lower().endswith('.pptx'):
                                from pptx import Presentation
                                prs = Presentation(file_path)
                                content = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
                            elif file_name.lower().endswith('.xlsx'):
                                import pandas as pd
                                content = ""
                                xl = pd.ExcelFile(file_path)
                                for sheet_name in xl.sheet_names:
                                    df = pd.read_excel(file_path, sheet_name=sheet_name)
                                    content += f"Sheet: {sheet_name}\n{df.to_string()}\n\n"
                        else:
                            # Skip unsupported file types
                            continue

                        # Get file metadata
                        stat = os.stat(file_path)
                        modified = datetime.datetime.fromtimestamp(stat.st_mtime).isoformat()

                        items.append({
                            "id": f"local_{hash(file_path)}",
                            "content": content,
                            "metadata": {
                                "source": "local",
                                "file_path": file_path,
                                "relative_path": relative_path,
                                "file_name": file_name,
                                "modified": modified,
                                "size": stat.st_size
                            }
                        })

                    except Exception as file_error:
                        logger.warning(f"Error processing file {file_path}: {str(file_error)}")
                        continue

        logger.info(f"Processed {len(items)} files from local directory")
        return items

    except Exception as e:
        logger.error(f"Error processing local source: {str(e)}")
        raise

async def process_item(item: Dict[str, Any], user_id: Optional[str] = None):
    """Process an individual item (extract text, generate embeddings, store in LanceDB)"""

    content = item.get("content", "")
    item_id = item.get("id", "")
    metadata = item.get("metadata", {})

    try:
        # 1. Clean and preprocess text content
        cleaned_content = clean_text_content(content)

        # 2. Generate embeddings using sentence transformers
        embeddings = await generate_embeddings(cleaned_content)

        # 3. Store in LanceDB vector database
        await store_in_lancedb(item_id, cleaned_content, embeddings, metadata, user_id)

        # 4. Index for search (optional - could be separate process)
        await index_for_search(item_id, cleaned_content, metadata)

        return {
            "status": "success",
            "item_id": item_id,
            "content_length": len(cleaned_content),
            "embedding_dimensions": len(embeddings) if embeddings else 0
        }

    except Exception as e:
        logger.error(f"Error processing item {item_id}: {str(e)}")
        return {
            "status": "error",
            "item_id": item_id,
            "error": str(e)
        }
    await asyncio.sleep(0.05)

    logger.debug(f"Processed item {item_id}: {content[:50]}...")

def clean_text_content(content: str) -> str:
    """Clean and preprocess text content"""
    import re
    import html

    # Decode HTML entities
    cleaned = html.unescape(content)

    # Remove extra whitespace and newlines
    cleaned = re.sub(r'\s+', ' ', cleaned)

    # Remove special characters but keep basic punctuation
    cleaned = re.sub(r'[^\w\s.,!?;:\-()\[\]{}]', '', cleaned)

    # Trim whitespace
    cleaned = cleaned.strip()

    return cleaned

async def generate_embeddings(text: str) -> List[float]:
    """Generate embeddings using sentence transformers"""
    from sentence_transformers import SentenceTransformer
    import numpy as np

    # Load model (cached for performance)
    if not hasattr(generate_embeddings, 'model'):
        generate_embeddings.model = SentenceTransformer('all-MiniLM-L6-v2')

    # Generate embeddings
    embedding = generate_embeddings.model.encode(text)

    return embedding.tolist()

async def store_in_lancedb(item_id: str, content: str, embedding: List[float], metadata: Dict[str, Any], user_id: Optional[str] = None):
    """Store item in LanceDB vector database"""
    import lancedb
    import pandas as pd

    # Connect to LanceDB
    db = lancedb.connect("/data/lancedb")

    # Create or get table
    table_name = "ingested_items"
    if table_name not in db.table_names():
        schema = lancedb.schema(
            [
                lancedb.field("id", lancedb.types.string()),
                lancedb.field("vector", lancedb.types.vector(len(embedding))),
                lancedb.field("content", lancedb.types.string()),
                lancedb.field("metadata", lancedb.types.json()),
                lancedb.field("user_id", lancedb.types.string()),
                lancedb.field("timestamp", lancedb.types.timestamp())
            ]
        )
        table = db.create_table(table_name, schema=schema)
    else:
        table = db.open_table(table_name)

    # Prepare data for insertion
    data = {
        "id": item_id,
        "vector": embedding,
        "content": content,
        "metadata": metadata,
        "user_id": user_id or "system",
        "timestamp": datetime.datetime.now()
    }

    # Insert into table
    table.add([data])

async def index_for_search(item_id: str, content: str, metadata: Dict[str, Any]):
    """Index content for search (optional)"""
    # This could integrate with Elasticsearch, Meilisearch, or other search engines
    # For now, we'll just log the indexing operation
    logger.info(f"Indexing item {item_id} for search: {content[:100]}...")

@app.post("/reset-counters")
async def reset_counters():
    """Reset processed items counter (for testing)"""
    global processed_items_count
    processed_items_count = 0
    ingestion_statuses.clear()
    return {"message": "Counters reset successfully"}

if __name__ == "__main__":
    uvicorn.run(app, host="0.0.0.0", port=8002)
